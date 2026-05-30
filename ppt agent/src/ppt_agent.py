from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt


def _chunk_text(text: str, max_len: int = 60) -> List[str]:
    """
    Very simple text splitter: splits by Chinese/English punctuation and newline,
    falls back to fixed-width chunks.
    """
    text = text.strip()
    if not text:
        return []

    # 优先按换行拆
    parts: List[str] = []
    for para in text.splitlines():
        para = para.strip()
        if not para:
            continue
        # 再按标点拆
        buf = ""
        for ch in para:
            buf += ch
            if ch in "。！？.!?" or len(buf) >= max_len:
                parts.append(buf.strip())
                buf = ""
        if buf.strip():
            parts.append(buf.strip())

    # 去掉过长的空行
    return [p for p in parts if p]


def generate_ppt(topic: str, content: str, output_path: str) -> Path:
    """
    生成一个最简单的 PPT：
    - 首页：标题 + 副标题（固定文案）
    - 后续 1~N 页：根据内容自动拆成要点，一页最多 5 条
    """
    prs = Presentation()

    # 标题页
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = topic
    subtitle.text = "由简单 PPT Agent 自动生成"

    # 内容页
    bullets = _chunk_text(content)
    if not bullets:
        bullets = [f"{topic} 的简介", "这里是自动生成的简单 PPT。"]

    max_per_slide = 5
    for i in range(0, len(bullets), max_per_slide):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        title.text = f"{topic}（第 {i // max_per_slide + 1} 部分）"

        tf = body.text_frame
        tf.clear()

        for j, bullet in enumerate(bullets[i : i + max_per_slide]):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

    output = Path(output_path).with_suffix(".pptx")
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


if __name__ == "__main__":
    # 简单自测入口：python -m src.ppt_agent "主题" "内容"
    import sys

    if len(sys.argv) < 3:
        print("用法: python -m src.ppt_agent <主题> <内容> [输出路径]")
        sys.exit(1)

    topic_arg = sys.argv[1]
    content_arg = sys.argv[2]
    out = sys.argv[3] if len(sys.argv) > 3 else f"{topic_arg}.pptx"

    path = generate_ppt(topic_arg, content_arg, out)
    print("已生成:", path)
